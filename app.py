import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import re
import requests
from PIL import Image
from copy import deepcopy

# --- File Paths – adjust as needed ---
MAPPING_FILE_PATH = "mapping-file.xlsx"
STOCK_FILE_PATH = "stock.xlsx"
TEMPLATE_FILE_PATH = "template-generator.pptx"

# --- Expected columns in the mapping file ---
REQUIRED_MAPPING_COLS_ORIG = [
    "{{Product name}}",
    "{{Product code}}",
    "{{Product country of origin}}",
    "{{Product height}}",
    "{{Product width}}",
    "{{Product length}}",
    "{{Product depth}}",
    "{{Product seat height}}",
    "{{Product diameter}}",
    "{{CertificateName}}",
    "{{Product Consumption COM}}",
    "{{Product Fact Sheet link}}",
    "{{Product configurator link}}",
    "{{Product Packshot1}}",
    "{{Product Lifestyle1}}",
    "{{Product Lifestyle2}}",
    "{{Product Lifestyle3}}",
    "{{Product Lifestyle4}}",
    "ProductKey"  # Mapping file's ProductKey (without brackets)
]

# --- Expected columns in the stock file ---
REQUIRED_STOCK_COLS_ORIG = [
    "productkey",    # Column B: ProductKey
    "variantname",   # Column D: VariantName
    "rts",           # Column H: RTS
    "mto"            # Column I: MTO
]

# --- Placeholders for text replacement in the template ---
TEXT_PLACEHOLDERS_ORIG = {
    "{{Product name}}": "Product Name:",
    "{{Product code}}": "Product Code:",
    "{{Product country of origin}}": "Country of origin:",
    "{{Product height}}": "Height:",
    "{{Product width}}": "Width:",
    "{{Product length}}": "Length:",
    "{{Product depth}}": "Depth:",
    "{{Product seat height}}": "Seat Height:",
    "{{Product diameter}}": "Diameter:",
    "{{CertificateName}}": "Test & certificates for the product:",
    "{{Product Consumption COM}}": "Consumption information for COM:"
}

# --- Placeholders for hyperlink replacement ---
HYPERLINK_PLACEHOLDERS_ORIG = {
    "{{Product Fact Sheet link}}": "Download Product Fact Sheet",
    "{{Product configurator link}}": "Click to configure product"
}

# --- Placeholders for image replacement ---
IMAGE_PLACEHOLDERS_ORIG = [
    "{{Product Packshot1}}",
    "{{Product Lifestyle1}}",
    "{{Product Lifestyle2}}",
    "{{Product Lifestyle3}}",
    "{{Product Lifestyle4}}",
]

# --- Helper Functions ---

def normalize_text(s):
    """Removes all whitespace (including non-breaking) and converts to lowercase."""
    return re.sub(r"\s+", "", str(s).replace("\u00A0", " ")).lower()

def normalize_col(col):
    """Normalizes column names for consistent matching."""
    return normalize_text(col)

def group_variant_names(variant_names, group_item_sep=", ", group_sep="\n"):
    """
    Groups a list of variant names based on the prefix (part before " - ").
    For each group, duplicates are removed, and the remaining parts (after " - ") are joined by group_item_sep.
    The groups are then joined by group_sep.
    """
    groups = {}
    for name in variant_names:
        if " - " in name:
            prefix, suffix = name.split(" - ", 1)
        else:
            prefix, suffix = name, ""
        prefix = prefix.strip()
        suffix = suffix.strip()
        groups.setdefault(prefix, set())
        if suffix:
            groups[prefix].add(suffix)
    output_lines = []
    for prefix, suffixes in sorted(groups.items()):
        suffix_list = sorted(list(suffixes))
        if suffix_list:
            line = f"{prefix} - " + group_item_sep.join(suffix_list)
        else:
            line = prefix
        output_lines.append(line)
    return group_sep.join(output_lines)

def find_mapping_row(item_no, mapping_df, mapping_prod_key):
    """
    Finds the row in mapping_df where the product code column (mapping_prod_key)
    matches 'Item no' (after normalization).
    """
    norm_item = normalize_text(item_no)
    # Exact match first
    for idx, row in mapping_df.iterrows():
        code = row.get(mapping_prod_key, "")
        if normalize_text(code) == norm_item:
            return row
    # Partial match if exact fails (e.g., '12345' for '12345-AB')
    if "-" in str(item_no):
        partial = normalize_text(str(item_no).split("-")[0])
        for idx, row in mapping_df.iterrows():
            code = row.get(mapping_prod_key, "")
            if normalize_text(code).startswith(partial):
                return row
    return None

def process_stock_rts_alternative(mapping_row, stock_df):
    """
    Logic for {{Product RTS}}:
    1. Get 'ProductKey' from mapping_row.
    2. Filter stock_df for matching 'productkey'.
    3. Filter again for non-empty 'rts'.
    4. Extract unique 'variantname' values.
    5. Group these values with group_variant_names().
    """
    product_key = mapping_row.get(normalize_col("ProductKey"), "")
    if not product_key or pd.isna(product_key):
        return ""
    norm_product_key = normalize_text(product_key)
    try:
        filtered = stock_df[stock_df["productkey"].apply(lambda x: normalize_text(x) == norm_product_key)]
    except KeyError as e:
        st.error(f"KeyError in RTS (productkey): {e}")
        return ""
    if filtered.empty:
        return ""
    filtered = filtered[filtered["rts"].notna() & (filtered["rts"] != "")]
    if filtered.empty:
        return ""
    try:
        variant_names = filtered["variantname"].dropna().astype(str).tolist()
    except KeyError as e:
        st.error(f"KeyError in RTS (variantname): {e}")
        return ""
    unique_variant_names = list(dict.fromkeys(variant_names))
    return group_variant_names(unique_variant_names, group_item_sep=", ", group_sep="\n")

def process_stock_mto_alternative(mapping_row, stock_df):
    """
    Logic for {{Product MTO}}:
    1. Get 'ProductKey' from mapping_row.
    2. Filter stock_df for matching 'productkey'.
    3. Filter again for non-empty 'mto'.
    4. Extract unique 'variantname' values.
    5. Group these values with group_variant_names().
    """
    product_key = mapping_row.get(normalize_col("ProductKey"), "")
    if not product_key or pd.isna(product_key):
        return ""
    norm_product_key = normalize_text(product_key)
    try:
        filtered = stock_df[stock_df["productkey"].apply(lambda x: normalize_text(x) == norm_product_key)]
    except KeyError as e:
        st.error(f"KeyError in MTO (productkey): {e}")
        return ""
    if filtered.empty:
        return ""
    filtered = filtered[filtered["mto"].notna() & (filtered["mto"] != "")]
    if filtered.empty:
        return ""
    try:
        variant_names = filtered["variantname"].dropna().astype(str).tolist()
    except KeyError as e:
        st.error(f"KeyError in MTO (variantname): {e}")
        return ""
    unique_variant_names = list(dict.fromkeys(variant_names))
    return group_variant_names(unique_variant_names, group_item_sep=", ", group_sep=", ")

def fetch_and_process_image(url, quality=70, max_size=(1200, 1200)):
    """Fetches, converts, resizes, and compresses an image from a URL."""
    if not isinstance(url, str) or not url.startswith('http'):
        return None
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status() # Raise an exception for bad status codes
        img_bytes = io.BytesIO(response.content)
        img = Image.open(img_bytes)
        
        # Convert to RGB to avoid issues with transparency (PNG, TIFF)
        if img.mode in ("RGBA", "LA", "P"):
            img = img.convert("RGB")
            
        img.thumbnail(max_size, Image.LANCZOS)
        
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format="JPEG", quality=quality, optimize=True)
        img_byte_arr.seek(0)
        return img_byte_arr
    except requests.exceptions.RequestException as e:
        st.warning(f"Network error fetching image from {url}: {e}")
    except Exception as e:
        st.warning(f"Could not process image from {url}: {e}")
    return None

def duplicate_slide(prs, slide):
    """Duplicates a slide and adds it to the end of the presentation."""
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Clear any default shapes from the new slide's layout
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy shapes from the source slide
    for shape in slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))
        
    return new_slide

def replace_text_in_shape(shape, placeholder_values):
    """Replaces text placeholders within a single shape."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        # Combine runs to find placeholders that might span across them
        full_text = "".join(run.text for run in paragraph.runs)
        if not any(ph in full_text for ph in placeholder_values):
            continue

        # Replace placeholders in the combined text
        new_text = full_text
        for placeholder, replacement in placeholder_values.items():
            # Ensure replacement is a string
            replacement_str = str(replacement) if replacement is not None else ""
            new_text = new_text.replace(placeholder, replacement_str)

        # Clear old runs and add a new one with the updated text
        if paragraph.runs:
            # Preserve formatting from the first run
            first_run = paragraph.runs[0]
            p = paragraph._p
            p.clear_content() # Removes all <a:r> child elements
            
            new_run = p.add_r()
            new_run.text = new_text
            # Copy font properties from the original first run
            if first_run.font:
                new_run.font.bold = first_run.font.bold
                new_run.font.italic = first_run.font.italic
                new_run.font.name = first_run.font.name
                new_run.font.size = first_run.font.size
                if first_run.font.color.type is not None:
                    new_run.font.color.rgb = first_run.font.color.rgb


def replace_hyperlink_placeholders(slide, hyperlink_values):
    """Replaces hyperlink placeholders in a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for placeholder, (display_text, url) in hyperlink_values.items():
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, display_text)
                        if url and isinstance(url, str):
                            try:
                                run.hyperlink.address = url
                            except Exception as e:
                                st.warning(f"Could not set hyperlink for {placeholder}: {e}")

def replace_image_placeholders(slide, image_values):
    """Replaces image placeholders with images from URLs."""
    shapes_to_remove = []
    pictures_to_add = []

    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text.strip():
            continue
        
        # Find a matching image placeholder in the shape's text
        for placeholder, url in image_values.items():
            if placeholder in shape.text:
                if url:
                    img_stream = fetch_and_process_image(url)
                    if img_stream:
                        pictures_to_add.append({
                            "stream": img_stream,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        })
                # Mark the placeholder shape for removal
                shapes_to_remove.append(shape)
                break # Move to the next shape once a placeholder is found

    # Add new pictures
    for pic in pictures_to_add:
        slide.shapes.add_picture(pic["stream"], pic["left"], pic["top"], width=pic["width"], height=pic["height"])
    
    # Remove old placeholder shapes
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)


def populate_slide_with_data(slide, mapping_row, stock_df, MAPPING_PRODUCT_CODE_KEY):
    """Helper function to populate a slide with data from a mapping row."""
    # --- Text Placeholders ---
    placeholder_texts = {}
    for ph, label in TEXT_PLACEHOLDERS_ORIG.items():
        norm_ph = normalize_col(ph)
        value = mapping_row.get(norm_ph, "")
        value = "" if pd.isna(value) else value
        
        if ph in ["{{CertificateName}}", "{{Product Consumption COM}}"]:
            placeholder_texts[ph] = f"{label}\n\n{value}"
        elif ph in ["{{Product name}}", "{{Product code}}", "{{Product country of origin}}"]:
             placeholder_texts[ph] = f"{label} {value}"
        else:
            placeholder_texts[ph] = f"{label}\n{value}"
    
    rts_text = process_stock_rts_alternative(mapping_row, stock_df)
    mto_text = process_stock_mto_alternative(mapping_row, stock_df)
    placeholder_texts["{{Product RTS}}"] = f"Product in stock versions:\n\n{rts_text}"
    placeholder_texts["{{Product MTO}}"] = f"Available for made to order:\n\n{mto_text}"

    for shape in slide.shapes:
        replace_text_in_shape(shape, placeholder_texts)

    # --- Hyperlink Placeholders ---
    hyperlink_vals = {}
    for ph, display_text in HYPERLINK_PLACEHOLDERS_ORIG.items():
        norm_ph = normalize_col(ph)
        url = mapping_row.get(norm_ph, "")
        url = "" if pd.isna(url) else url
        hyperlink_vals[ph] = (display_text, url)
    replace_hyperlink_placeholders(slide, hyperlink_vals)

    # --- Image Placeholders ---
    image_vals = {}
    for ph in IMAGE_PLACEHOLDERS_ORIG:
        norm_ph = normalize_col(ph)
        url = mapping_row.get(norm_ph, "")
        url = "" if pd.isna(url) else url
        image_vals[ph] = url
    replace_image_placeholders(slide, image_vals)


# --- Main App ---
def main():
    st.title("PowerPoint Generator App")
    st.write("Enter item numbers, one per line:")
    st.info("Note: Enter item numbers without extra spaces around hyphens, e.g., '03084' or '12345-AB'.")
    pasted_text = st.text_area("Paste item numbers here", height=200)
    
    if not pasted_text.strip():
        st.info("Please enter item numbers in the text area.")
        return

    item_numbers = [line.strip() for line in pasted_text.splitlines() if line.strip()]
    if not item_numbers:
        st.error("No valid item numbers found.")
        return

    user_df = pd.DataFrame({"Item no": item_numbers})

    st.write("User data created successfully!")
    st.info("Validating files...")
    progress_bar = st.progress(10)

    # Load mapping file
    try:
        mapping_df = pd.read_excel(MAPPING_FILE_PATH)
        mapping_df.columns = [normalize_col(col) for col in mapping_df.columns]
    except Exception as e:
        st.error(f"Error reading mapping file: {e}")
        return

    normalized_required_mapping_cols = [normalize_col(col) for col in REQUIRED_MAPPING_COLS_ORIG]
    missing_mapping_cols = [req for req in normalized_required_mapping_cols if req not in mapping_df.columns]
    if missing_mapping_cols:
        st.error(f"Mapping file is missing columns (after normalization): {missing_mapping_cols}.")
        return

    st.write("Mapping file loaded successfully!")
    progress_bar.progress(30)
    MAPPING_PRODUCT_CODE_KEY = normalize_col("{{Product code}}")

    # Load stock file
    try:
        stock_df = pd.read_excel(STOCK_FILE_PATH)
        stock_df.columns = [normalize_col(col) for col in stock_df.columns]
    except Exception as e:
        st.error(f"Error reading stock file: {e}")
        return

    normalized_required_stock_cols = [normalize_col(col) for col in REQUIRED_STOCK_COLS_ORIG]
    missing_stock_cols = [req for req in normalized_required_stock_cols if req not in stock_df.columns]
    if missing_stock_cols:
        st.error(f"Stock file is missing columns (after normalization): {missing_stock_cols}.")
        return

    st.write("Stock file loaded successfully!")
    progress_bar.progress(50)

    # Load PowerPoint template
    try:
        prs = Presentation(TEMPLATE_FILE_PATH)
    except Exception as e:
        st.error(f"Error reading template file: {e}")
        return

    if len(prs.slides) < 1:
        st.error("Template file must contain at least one slide.")
        return

    st.write("Template file loaded successfully!")
    progress_bar.progress(70)

    # --- CORE LOGIC: Use first slide for first item, then duplicate ---
    template_slide = prs.slides[0]
    
    # Process the first product on the original template slide
    first_product = user_df.iloc[0]
    item_no = first_product["Item no"]
    mapping_row = find_mapping_row(item_no, mapping_df, MAPPING_PRODUCT_CODE_KEY)
    
    if mapping_row is not None:
        populate_slide_with_data(template_slide, mapping_row, stock_df, MAPPING_PRODUCT_CODE_KEY)
    else:
        st.warning(f"No match found for the first item: {item_no}. First slide will be blank.")
        # Clear all placeholders on the first slide if no match
        for shape in template_slide.shapes:
            if shape.has_text_frame:
                shape.text = ""

    progress_bar.progress(70 + int(30 / len(user_df)))

    # Loop through the rest of the products and duplicate the slide
    for index, product in user_df.iloc[1:].iterrows():
        item_no = product["Item no"]
        
        # Create a new slide by duplicating the original template slide
        new_slide = duplicate_slide(prs, template_slide)

        mapping_row = find_mapping_row(item_no, mapping_df, MAPPING_PRODUCT_CODE_KEY)
        if mapping_row is None:
            st.warning(f"No match found for item: {item_no}. A blank slide will be added.")
            # Clear all placeholders on the new slide
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    shape.text = ""
            continue

        populate_slide_with_data(new_slide, mapping_row, stock_df, MAPPING_PRODUCT_CODE_KEY)

        progress_value = 70 + min(int((index + 1) / len(user_df) * 30), 30)
        progress_bar.progress(progress_value)

    # --- Save and Download ---
    ppt_io = io.BytesIO()
    try:
        prs.save(ppt_io)
        ppt_io.seek(0)
    except Exception as e:
        st.error(f"Error saving PowerPoint: {e}")
        return

    st.success("PowerPoint generated successfully!")
    st.download_button("Download PowerPoint", ppt_io,
                       file_name="generated_presentation.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == '__main__':
    main()
